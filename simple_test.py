"""
Simple test without MCP protocol
"""
import os
from dotenv import load_dotenv

# Load environment variables
load_dotenv()

print("Testing S3 Configuration...")
print(f"S3_ENABLED: {os.getenv('S3_ENABLED')}")
print(f"S3_BUCKET_NAME: {os.getenv('S3_BUCKET_NAME')}")
print(f"S3_PREFIX: {os.getenv('S3_PREFIX')}")
print(f"AWS_DEFAULT_REGION: {os.getenv('AWS_DEFAULT_REGION')}")

# Test importing S3 utils
try:
    from utils.s3_utils import get_s3_handler
    print("\n✅ Successfully imported S3 utils")
    
    s3_handler = get_s3_handler()
    print(f"✅ S3 Handler initialized")
    print(f"   Enabled: {s3_handler.s3_enabled}")
    
    if s3_handler.s3_enabled:
        print(f"   Bucket: {s3_handler.s3_bucket}")
        print(f"   Prefix: {s3_handler.s3_prefix}")
        print(f"   Region: {s3_handler.s3_region}")
        
        # Test creating and uploading a presentation
        print("\n📄 Creating test presentation...")
        from pptx import Presentation
        
        prs = Presentation()
        title_slide = prs.slide_layouts[0]
        slide = prs.slides.add_slide(title_slide)
        title = slide.shapes.title
        title.text = "Test from Simple Script"
        
        print("📤 Uploading to S3...")
        result = s3_handler.upload_presentation(prs, "simple_test.pptx")
        
        if result["success"]:
            print(f"\n✅ SUCCESS!")
            print(f"   S3 URL: {result['s3_url']}")
            print(f"   HTTPS URL: {result['https_url']}")
        else:
            print(f"\n❌ FAILED: {result.get('error')}")
    
except Exception as e:
    print(f"\n❌ Error: {e}")
    import traceback
    traceback.print_exc()